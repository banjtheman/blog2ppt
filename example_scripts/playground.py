import json
import pandas as pd
import logging
import utils
from pptx import Presentation
from pptx.util import Inches
import urllib.request
import carbon_img_gen
import asyncio


def create_ppt(title: str, desc: str, cover_img: str):
    print("Creating pttx")

    # Creating presentation object
    ppt = Presentation()

    # Creating slide layout
    first_slide_layout = ppt.slide_layouts[0]

    """ Ref for slide types: 
    0 ->  title and subtitle
    1 ->  title and content
    2 ->  section header
    3 ->  two content
    4 ->  Comparison
    5 ->  Title only 
    6 ->  Blank
    7 ->  Content with caption
    8 ->  Pic with caption
    """

    ############ Cover slide ###########
    # Creating slide object to add
    # in ppt i.e. Attaching slides
    # with Presentation i.e. ppt
    slide = ppt.slides.add_slide(first_slide_layout)

    # Adding title and subtitle in
    # slide i.e. first page of slide
    slide.shapes.title.text = title

    # We have different formats of
    # subtitles in ppts, for simple
    # subtitle this method should
    # implemented, you can change
    # 0 to 1 for different design
    slide.placeholders[1].text = desc

    # For margins
    left = top = Inches(1)

    # adding image
    # slide.shapes.add_picture(cover_img, left, top)

    # Add smaller picture
    left = Inches(1)
    height = Inches(1)

    slide.shapes.add_picture(cover_img, left, top, height=height)
    ############ Cover slide ###########

    # save file
    # ppt.save('test_4.pptx')
    return ppt


def get_blog_details(text):
    """
    Purpose:
        Get details of the blog
    Args:
        N/A
    Returns:
        N/A
    """
    blog_json = {}

    lines = text.split("\n")
    # title
    title = lines[1].split(": ")[1]
    # description
    desc = lines[3].split(": ")[1]
    # Cover image
    cover_image = lines[5].split(": ")[1]

    # Download image
    image_name = cover_image.split("/")[-1]
    image_dl_loc = f"downloaded_images/{image_name}"
    urllib.request.urlretrieve(cover_image, image_dl_loc)

    blog_json["title"] = title
    blog_json["desc"] = desc
    blog_json["cover_image"] = image_dl_loc
    blog_json["sections"] = []


    # Create sections based on Headings (hmm inital heading counts as well)
    #  for each Section
    #   find code snippets (convert to carbon with carbon mod)
    #   find images (download imae with urllib)  
    #   summarize text (using transformer model)
    #   Create Sepatartor slide (if init section no need for sepator slide)
    #   Create slides per sentence in summary (using spacy to detect) 
    #   Create Slides for each image/code snippet



    return blog_json


def main() -> None:
    """
    Purpose:
        Main script
    Args:
        N/A
    Returns:
        N/A
    """
    print("This is the way")

    # We have mechanism to transform code into carbon.io img
    # code_img = asyncio.run(carbon_img_gen.get_carbon_image(code))
    # return

    # Read text file
    curr_blog_post = utils.read_from_file("example.txt")

    # Get metadata
    # Get each line
    lines = curr_blog_post.split("\n")
    # title
    title = lines[1].split(": ")[1]
    # description
    desc = lines[3].split(": ")[1]
    # Cover image
    cover_image = lines[5].split(": ")[1]

    # Download image
    image_name = cover_image.split("/")[-1]
    image_dl_loc = f"downloaded_images/{image_name}"
    urllib.request.urlretrieve(cover_image, image_dl_loc)

    # Create a PPT
    ppt = create_ppt(title, desc, image_dl_loc)

    # Save ppt
    ppt.save("ppts/test_4.pptx")

    print("Done and done")


if __name__ == "__main__":
    loglevel = logging.INFO
    logging.basicConfig(format="%(levelname)s: %(message)s", level=loglevel)
    main()
