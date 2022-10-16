# Python imports
import requests
from typing import Type, Union, Dict, Any, List
import os
os.environ['TRANSFORMERS_CACHE'] = '/tmp/' # Set tmp as transformers cache

# 3rd party imports
from bs4 import BeautifulSoup
import urllib.request
from pptx import Presentation
from pptx.util import Inches
from transformers import pipeline
import asyncio
from carbon import Carbon
from spacy.lang.en import English
import math

# My imports
import utils


# Constants
carbon_client = Carbon()
spacy_nlp = English()  # Find sentences
spacy_nlp.add_pipe("sentencizer")


title_bait = "- DEV Community 👩‍💻👨‍💻"  # Remove from title

# Huggingface summarizer
summarizer = pipeline(
    "summarization", model="facebook/bart-large-cnn"
)  # Text summarizer


def split_in_sentences(text):
    """
    Purpose:
        Split up sentences
    Args/Requests:
         text: text to split
    Return:
        array with sentences
    """
    doc = spacy_nlp(text)
    return [str(sent).strip() for sent in doc.sents]


async def get_carbon_image(code: str) -> str:
    """
    Purpose:
        Create carbon file
    Args/Requests:
         code: code to carbonize
    Return:
        file_path: loc of carbon file
    """
    img = await carbon_client.create(code, downloads_dir="/tmp/")
    return img


def create_new_section(name: str):
    """
    Purpose:
       Create new section
    Args:
        name - name of section
    Returns:
        section json
    """

    curr_section = {}
    # curr_section["text"] = []
    # curr_section["images"] = []
    # curr_section["images_alt"] = []
    # curr_section["code"] = []
    # curr_section["order"] = []
    # curr_section["lists"] = []

    curr_section["name"] = name
    curr_section["elements"] = []
    # TODO add links and last slide can be all the stuff you linked throughout post

    return curr_section


def is_any_heading(text: str):
    """
    Purpose:
       Check if tag is a heading
    Args:
        text - text to check
    Returns:
        true if heading, else false
    """
    if text == "h1":
        return True
    elif text == "h2":
        return True
    elif text == "h3":
        return True
    elif text == "h4":
        return True
    elif text == "h5":
        return True
    elif text == "h6":
        return True
    else:
        return False


def get_html_sections(url: str):
    """
    Purpose:
       Get article text
    Args:
        url - the url
    Returns:
        text of the article
    """
    raw_html = requests.get(url, timeout=30).text
    soup = BeautifulSoup(raw_html, "html.parser")
    # article_text = soup.find(id="article-body").text  # Raw text of article

    blog_json = {}
    num_sections = 0

    # title
    title = soup.find("title").string.replace(title_bait, "").strip()
    # description
    desc_tag = soup.find("meta", property="og:description")
    if desc_tag:
        desc = desc_tag.get("content")
    else:
        desc = ""
    # Cover image
    cover_tag = soup.find("meta", property="og:image")
    if cover_tag:
        cover_text = cover_tag.get("content")
        cover_image = f'https://{cover_text.split("https://")[2]}'

    # Download image
    try:
        image_name = cover_image.split("/")[-1]
        image_dl_loc = f"/tmp/downloaded_images/{image_name}"
        # Just to be safe
        os.system("mkdir -p /tmp/downloaded_images/")
        urllib.request.urlretrieve(cover_image, image_dl_loc)
    except:
        print("Url not found...")
        print(cover_image)
        try:
            image_name = "cover_image.png"
            image_dl_loc = f"/tmp/downloaded_images/{image_name}"
            urllib.request.urlretrieve(cover_text, image_dl_loc)
        except:
            print("Can't get image...")
            print(cover_text)
            cover_image = None
            image_dl_loc = None


    blog_json["title"] = title
    blog_json["desc"] = desc
    blog_json["cover_image"] = image_dl_loc
    blog_json["sections"] = []

    main_article = soup.find(id="article-body")  # Raw text of article

    # Base section
    curr_section = create_new_section("Intro")

    # Classes for code blocks
    code_classes = ["highlight", "js-code-highlight"]

    for element in main_article.children:

        # Based on child type, do certain actions

        if element.name == "ul" or element.name == "ol":

            curr_list = []
            list_items = element.findAll("li")

            for list_item in list_items:
                curr_list.append(list_item.text)

            curr_elem = {"elem_type": "list", "data": curr_list}
            curr_section["elements"].append(curr_elem)
            num_sections += 1

        # Check for div with code
        elif element.name == "div":
            # check for code block
            if element["class"] == code_classes:

                # Get all code blocks
                curr_code = element.findAll("code")
                # Add code
                for code_snippets in curr_code:

                    curr_elem = {"elem_type": "code", "data": code_snippets.text}
                    curr_section["elements"].append(curr_elem)
                    num_sections += 1

        # if element is p
        elif element.name == "p":

            # make sure its actual text
            if len(element.text) > 0:

                curr_elem = {"elem_type": "text", "data": element.text}
                curr_section["elements"].append(curr_elem)

                num_sections += 1

            # Check if any images
            curr_images = element.findAll("img")

            for image in curr_images:
                image_url = image["src"]
                true_image_url = f'https://{image_url.split("https://")[2]}'
                image_text = image["alt"]

                curr_elem = {
                    "elem_type": "image",
                    "data": true_image_url,
                    "alt_text": image_text,
                }
                curr_section["elements"].append(curr_elem)
                num_sections += 1

        # Check for any of the headings
        elif is_any_heading(element.name):
            # Starting a new section
            blog_json["sections"].append(curr_section)

            # Reset current_section
            clean_name = element.text.replace("\n", "").strip()
            curr_section = create_new_section(clean_name)

            # num_sections += 1

        # print(element.name)
        # print(element)
        # print("##############NEXT##########")

    # add last section
    blog_json["sections"].append(curr_section)

    # print(blog_json)

    return blog_json, num_sections


async def convert_sections_to_ppt(blog_sections, num_sections, prog_bar=None):
    """
    Purpose:
       convert blog sections into powerpoint
    Args:
        blog_sections - the sections of the blog
    Returns:
        The powerpoint object
    """
    print("Creating pttx")

    # Creating presentation object
    ppt = Presentation()

    # Creating slide layout
    first_slide_layout = ppt.slide_layouts[0]
    bullet_slide_layout = ppt.slide_layouts[1]
    section_layout = ppt.slide_layouts[2]
    blank_layout = ppt.slide_layouts[6]
    pic_layout = ppt.slide_layouts[8]

    """ Ref for slide types: 
    0 ->  title and subtitle
    1 ->  title and content
    2 ->  section header
    3 ->  two content
    4 ->  Comparison
    5 ->  Title only 
    6 ->  Blank
    7 ->  Content with caption
    8 ->  Pic with caption
    """
    # Cover slide info
    title = blog_sections["title"]
    desc = blog_sections["desc"]
    cover_image = blog_sections["cover_image"]

    # Margins
    top = Inches(1)
    left = Inches(1)
    # left = int((ppt.slide_width - image.width) / 2)
    height = Inches(4)

    # Make picture slide
    slide = ppt.slides.add_slide(pic_layout)

    placeholder = slide.placeholders[1]
    slide.shapes.title.text = title
    if cover_image:
        placeholder.insert_picture(cover_image)
    slide.placeholders[2].text = desc

    # prog_update = int(100 / num_sections)
    # prog_update = math.ceil(100 / num_sections)
    prog_update = 1 / num_sections

    curr_prog = 0
    print(f"Num sections: {num_sections}")

    # Create slides for each section
    for index, section in enumerate(blog_sections["sections"]):

        section_name = section["name"]
        # check if not Intro section

        if index != 0:
            print("Create section header")
            slide = ppt.slides.add_slide(section_layout)
            slide.shapes.title.text = section_name

        # Now create slides for each part in section
        for element in section["elements"]:

            print(f"Working on element: {element}")

            elem_type = element["elem_type"]

            if elem_type == "text":
                # Summarize text
                curr_text = element["data"]
                # curr_len = len(curr_text)
                summed_text = summarizer(
                    curr_text, max_length=80, min_length=30, do_sample=False
                )

                slide_text = summed_text[0]["summary_text"]
                curr_sentences = split_in_sentences(slide_text)

                slide = ppt.slides.add_slide(bullet_slide_layout)
                shapes = slide.shapes

                # Section title
                title_shape = shapes.title
                title_shape.text = section_name

                body_shape = shapes.placeholders[1]

                for index, sent in enumerate(curr_sentences):
                    # Init bullet points
                    if index == 0:
                        text_frame = body_shape.text_frame
                        text_frame.text = sent
                    # Write new line
                    else:
                        next_line = text_frame.add_paragraph()
                        next_line.text = sent

            elif elem_type == "list":

                slide = ppt.slides.add_slide(bullet_slide_layout)
                shapes = slide.shapes

                # Section title
                title_shape = shapes.title
                title_shape.text = section_name

                body_shape = shapes.placeholders[1]

                for index, list_ele in enumerate(element["data"]):
                    # Init bullet points
                    if index == 0:
                        text_frame = body_shape.text_frame
                        text_frame.text = list_ele
                    # Write new line
                    else:
                        next_line = text_frame.add_paragraph()
                        next_line.text = list_ele

            elif elem_type == "code":
                curr_code = element["data"]
                img_path = await carbon_client.create(curr_code, downloads_dir="/tmp/")

                # Create slide with code snippet
                slide = ppt.slides.add_slide(blank_layout)
                slide.shapes.add_picture(img_path, left, top, height=height)

                # slide = ppt.slides.add_slide(pic_layout)
                # placeholder = slide.placeholders[1]
                # placeholder.insert_picture(img_path)

            elif elem_type == "image":
                curr_image = element["data"]

                image_name = curr_image.split("/")[-1]
                image_dl_loc = f"/tmp/downloaded_images/{image_name}"
                urllib.request.urlretrieve(curr_image, image_dl_loc)

                # Create slide with image
                slide = ppt.slides.add_slide(blank_layout)
                slide.shapes.add_picture(image_dl_loc, left, top, height=height)

                # slide = ppt.slides.add_slide(pic_layout)
                # placeholder = slide.placeholders[1]
                # placeholder.insert_picture(image_dl_loc)
            else:
                print("Invlaid element????")

            if prog_bar:
                curr_prog += prog_update
                print(f"Current progress: {curr_prog}")
                if curr_prog > 1.0:
                    curr_prog = 0.99
                prog_bar.progress(curr_prog)

    if prog_bar:
        prog_bar.progress(1.0)
    return ppt


async def main():
    print("loading data...")
    blog_url = (
        "https://dev.to/banjtheman/dc-fire-and-ems-data-visualizer-4a6l"
    )

    print(f"Converting blog {blog_url} to ppt")

    blog_sections, num_sections = get_html_sections(blog_url)
    # utils.save_json("example.json", blog_sections)

    # utils.write_to_file("example.html", html_text)

    ppt = await convert_sections_to_ppt(blog_sections, num_sections)
    ppt.save("end_to_end.pptx")

    print("done and done")


if __name__ == "__main__":
    asyncio.run(main())
    # main()
